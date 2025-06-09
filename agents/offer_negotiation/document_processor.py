"""Document processing module for handling various document formats."""

import logging
from pathlib import Path
from typing import Optional, Union

import pdfplumber
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

# Supported document extensions
SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF Document",
    ".doc": "Word Document",
    ".docx": "Word Document",
    ".ppt": "PowerPoint Presentation",
    ".pptx": "PowerPoint Presentation",
}


def is_supported_document(file_path: Union[str, Path]) -> bool:
    """
    Check if the file is a supported document type.

    Args:
        file_path: Path to the file to check

    Returns:
        bool: True if the file is a supported document type
    """
    file_path = Path(file_path)
    return file_path.suffix.lower() in SUPPORTED_EXTENSIONS


def extract_text_from_pdf(file_path: Path) -> Optional[str]:
    """
    Extract text from a PDF file using pdfplumber.

    Args:
        file_path: Path to the PDF file

    Returns:
        str: Extracted text from the PDF, or None if extraction failed
    """
    try:
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text() or "")
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
        return None


def extract_text_from_docx(file_path: Path) -> Optional[str]:
    """
    Extract text from a Word document.

    Args:
        file_path: Path to the Word document

    Returns:
        str: Extracted text from the document, or None if extraction failed
    """
    try:
        doc = Document(file_path)
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as e:
        logger.error(f"Error extracting text from Word document {file_path}: {str(e)}")
        return None


def extract_text_from_pptx(file_path: Path) -> Optional[str]:
    """
    Extract text from a PowerPoint presentation.

    Args:
        file_path: Path to the PowerPoint presentation

    Returns:
        str: Extracted text from the presentation, or None if extraction failed
    """
    try:
        prs = Presentation(file_path)
        text = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text.append("\n".join(slide_text))

        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint {file_path}: {str(e)}")
        return None


def extract_text_from_document(file_path: Union[str, Path]) -> Optional[str]:
    """
    Extract text from a document file based on its extension.

    Args:
        file_path: Path to the document file

    Returns:
        str: Extracted text from the document, or None if extraction failed
    """
    file_path = Path(file_path)

    if not file_path.exists():
        logger.error(f"File not found: {file_path}")
        return None

    if not is_supported_document(file_path):
        logger.error(f"Unsupported document type: {file_path.suffix}")
        return None

    extension = file_path.suffix.lower()

    try:
        if extension == ".pdf":
            return extract_text_from_pdf(file_path)
        elif extension in [".doc", ".docx"]:
            return extract_text_from_docx(file_path)
        elif extension in [".ppt", ".pptx"]:
            return extract_text_from_pptx(file_path)
        else:
            logger.error(f"Unsupported document type: {extension}")
            return None
    except Exception as e:
        logger.error(f"Error processing document {file_path}: {str(e)}")
        return None
